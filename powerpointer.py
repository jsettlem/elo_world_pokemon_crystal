import colorsys
import json
import pickle
import zlib
from pprint import pprint
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor, ColorFormat
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from tqdm import tqdm

from analysis_models.trainer import Trainer
from protobuf.game_data_pb2 import MoveIdentifier
from utils.data import moves
from utils.hashids import prettify_hashid

with open("raw_trainer_movesets.json", "r") as f:
    raw_trainer_movesets = json.loads(f.read())

with open("./data_files/species_types.json", "r") as f:
    species_types = json.loads(f.read())

type_colors = {
    "NORMAL": RGBColor.from_string("a8a878"),
    "FIRE": RGBColor.from_string("f08030"),
    "WATER": RGBColor.from_string("6890f0"),
    "ELECTRIC": RGBColor.from_string("f8d030"),
    "GRASS": RGBColor.from_string("78c850"),
    "ICE": RGBColor.from_string("98d8d8"),
    "FIGHTING": RGBColor.from_string("c03028"),
    "POISON": RGBColor.from_string("a040a0"),
    "GROUND": RGBColor.from_string("e0c068"),
    "FLYING": RGBColor.from_string("a890f0"),
    "PSYCHIC_TYPE": RGBColor.from_string("f85888"),
    "BUG": RGBColor.from_string("a8b820"),
    "ROCK": RGBColor.from_string("b8a038"),
    "GHOST": RGBColor.from_string("705898"),
    "DRAGON": RGBColor.from_string("7038f8"),
    "DARK": RGBColor.from_string("705848"),
    "STEEL": RGBColor.from_string("b8b8d0"),
    "FAIRY": RGBColor.from_string("ee99ac"),
    "CURSE_TYPE": RGBColor.from_string("68A090"),
    "NONE": RGBColor.from_string("ffffff"),
}

pokemon_in_order = ["Bulbasaur", "Ivysaur", "Venusaur", "Charmander", "Charmeleon", "Charizard", "Squirtle",
                    "Wartortle", "Blastoise", "Caterpie", "Metapod", "Butterfree", "Weedle", "Kakuna",
                    "Beedrill", "Pidgey", "Pidgeotto", "Pidgeot", "Rattata", "Raticate", "Spearow", "Fearow",
                    "Ekans", "Arbok", "Pikachu", "Raichu", "Sandshrew", "Sandslash", "Nidoran♀", "Nidorina",
                    "Nidoqueen", "Nidoran♂", "Nidorino", "Nidoking", "Clefairy", "Clefable", "Vulpix",
                    "Ninetales", "Jigglypuff", "Wigglytuff", "Zubat", "Golbat", "Oddish", "Gloom", "Vileplume",
                    "Paras", "Parasect", "Venonat", "Venomoth", "Diglett", "Dugtrio", "Meowth", "Persian",
                    "Psyduck", "Golduck", "Mankey", "Primeape", "Growlithe", "Arcanine", "Poliwag", "Poliwhirl",
                    "Poliwrath", "Abra", "Kadabra", "Alakazam", "Machop", "Machoke", "Machamp", "Bellsprout",
                    "Weepinbell", "Victreebel", "Tentacool", "Tentacruel", "Geodude", "Graveler", "Golem",
                    "Ponyta", "Rapidash", "Slowpoke", "Slowbro", "Magnemite", "Magneton", "Farfetch'd", "Doduo",
                    "Dodrio", "Seel", "Dewgong", "Grimer", "Muk", "Shellder", "Cloyster", "Gastly", "Haunter",
                    "Gengar", "Onix", "Drowzee", "Hypno", "Krabby", "Kingler", "Voltorb", "Electrode",
                    "Exeggcute", "Exeggutor", "Cubone", "Marowak", "Hitmonlee", "Hitmonchan", "Lickitung",
                    "Koffing", "Weezing", "Rhyhorn", "Rhydon", "Chansey", "Tangela", "Kangaskhan", "Horsea",
                    "Seadra", "Goldeen", "Seaking", "Staryu", "Starmie", "Mr. Mime", "Scyther", "Jynx",
                    "Electabuzz", "Magmar", "Pinsir", "Tauros", "Magikarp", "Gyarados", "Lapras", "Ditto",
                    "Eevee", "Vaporeon", "Jolteon", "Flareon", "Porygon", "Omanyte", "Omastar", "Kabuto",
                    "Kabutops", "Aerodactyl", "Snorlax", "Articuno", "Zapdos", "Moltres", "Dratini",
                    "Dragonair", "Dragonite", "Mewtwo", "Mew",
                    "Chikorita", "Bayleef", "Meganium", "Cyndaquil", "Quilava", "Typhlosion", "Totodile",
                    "Croconaw", "Feraligatr", "Sentret", "Furret", "Hoothoot", "Noctowl", "Ledyba", "Ledian",
                    "Spinarak", "Ariados", "Crobat", "Chinchou", "Lanturn", "Pichu", "Cleffa", "Igglybuff",
                    "Togepi", "Togetic", "Natu", "Xatu", "Mareep", "Flaaffy", "Ampharos", "Bellossom",
                    "Marill", "Azumarill", "Sudowoodo", "Politoed", "Hoppip", "Skiploom", "Jumpluff",
                    "Aipom", "Sunkern", "Sunflora", "Yanma", "Wooper", "Quagsire", "Espeon", "Umbreon",
                    "Murkrow", "Slowking", "Misdreavus", "Unown", "Wobbuffet", "Girafarig", "Pineco",
                    "Forretress", "Dunsparce", "Gligar", "Steelix", "Snubbull", "Granbull", "Qwilfish",
                    "Scizor", "Shuckle", "Heracross", "Sneasel", "Teddiursa", "Ursaring", "Slugma", "Magcargo",
                    "Swinub", "Piloswine", "Corsola", "Remoraid", "Octillery", "Delibird", "Mantine",
                    "Skarmory", "Houndour", "Houndoom", "Kingdra", "Phanpy", "Donphan", "Porygon2",
                    "Stantler", "Smeargle", "Tyrogue", "Hitmontop", "Smoochum", "Elekid", "Magby", "Miltank",
                    "Blissey", "Raikou", "Entei", "Suicune", "Larvitar", "Pupitar", "Tyranitar", "Lugia",
                    "Ho-Oh", "Celebi"
                    ]

item_sprites = {
    'FOCUS_BAND': "FOCUS_BAND.png",
    'HYPER_POTION': "Bag_Hyper_Potion_Sprite.png",
    'KINGS_ROCK': "Bag_King's_Rock_Sprite.png",
    'DIRE_HIT': "Bag_Dire_Hit_Sprite.png",
    'PINK_BOW': "PINK_BOW.png",
    'MAX_POTION': "Bag_Max_Potion_Sprite.png",
    'FULL_RESTORE': "Bag_Full_Restore_Sprite.png",
    'FULL_HEAL': "Bag_Full_Heal_Sprite.png",
    'BERRY': "Bag_Oran_Berry_Sprite.png",
}

trainer_sprites = {
    59: "Spr_GS_Pokéfan_M.png",
    30: "Spr_GS_PokéManiac.png",
    29: "Spr_GS_Beauty_JP.png",
    45: "Spr_GS_Biker.png",
    24: "Spr_GS_Bird_Keeper.png",
    50: "Spr_GS_Blackbelt.png",
    58: "Spr_GS_Boarder.png",
    36: "Spr_GS_Bug_Catcher.png",
    47: "Spr_GS_Burglar.png",
    54: "Spr_GS_Camper.png",
    16: "Spr_GS_Lance.png",
    28: "Spr_GS_Cooltrainer_F.png",
    27: "Spr_GS_Cooltrainer_M.png",
    13: "Spr_GS_Bruno.png",
    14: "Spr_GS_Karen.png",
    15: "Spr_GS_Koga.png",
    11: "Spr_GS_Will.png",
    48: "Spr_GS_Firebreather.png",
    37: "Spr_GS_Fisher.png",
    32: "Spr_GS_Gentleman.png",
    43: "Spr_GS_Guitarist.png",
    44: "Spr_GS_Hiker.png",
    49: "Spr_GS_Juggler.png",
    60: "Spr_GS_Kimono_Girl.png",
    25: "Spr_GS_Lass.png",
    46: "Spr_GS_Blaine.png",
    64: "Spr_GS_Blue.png",
    17: "Spr_GS_Brock.png",
    3: "Spr_GS_Bugsy.png",
    7: "Spr_GS_Chuck.png",
    8: "Spr_GS_Clair.png",
    21: "Spr_GS_Erika.png",
    1: "Spr_GS_Falkner.png",
    26: "Spr_GS_Janine.png",
    6: "Spr_GS_Jasmine.png",
    19: "Spr_GS_Lt_Surge.png",
    18: "Spr_GS_Misty.png",
    4: "Spr_GS_Morty.png",
    5: "Spr_GS_Pryce.png",
    35: "Spr_GS_Sabrina.png",
    2: "Spr_GS_Whitney.png",
    57: "Spr_GS_Medium_JP.png",
    67: "Spr_C_Eusine.png",
    65: "Spr_GS_Officer.png",
    53: "Spr_GS_Picnicker.png",
    12: "Spr_C_Ethan.png",
    63: "Spr_GS_Red.png",
    62: "Spr_GS_Pokéfan_F.png",
    52: "Spr_GS_Psychic.png",
    9: "Spr_GS_Silver_1.png",
    42: "Spr_GS_Silver_2.png",
    51: "Spr_GS_Rocket_Executive_M.png",
    55: "Spr_GS_Rocket_Executive_F.png",
    31: "Spr_GS_Rocket_Grunt_M.png",
    66: "Spr_GS_Rocket_Grunt_F.png",
    56: "Spr_GS_Sage_JP.png",
    40: "Spr_GS_Sailor.png",
    23: "Spr_GS_Schoolboy.png",
    20: "Spr_GS_Scientist.png",
    33: "Spr_GS_Skier.png",
    41: "Spr_GS_Super_Nerd.png",
    39: "Spr_GS_Swimmer_F_JP.png",
    38: "Spr_GS_Swimmer_M.png",
    34: "Spr_GS_Teacher.png",
    61: "Spr_GS_Twins.png",
    22: "Spr_GS_Youngster.png",
}


def format_name(name: str):
    match name:
        case "?":
            return "???"
        case "#MANIAC":
            return "PokéManiac"
        case "#FAN" | "POKéFAN":
            return "Pokéfan"
        case "PKMN TRAINER":
            return "Pokémon Trainer"
        case "SWIMMER♂" | "SWIMMER♀":
            return "Swimmer"
        case "NIDORAN_F":
            return "Nidoran♀"
        case "NIDORAN_M":
            return "Nidoran♂"
        case "MR__MIME":
            return "Mr. Mime"
        case "FARFETCH_D":
            return "Farfetch'd"
        case "PSYCHIC_M":
            return "Psychic"
        case _:
            return name.replace("_", " ").title()


def getSuffix(n):
    if n % 100 in [11, 12, 13]:
        return 'th'
    match n % 10:
        case 1:
            return 'st'
        case 2:
            return 'nd'
        case 3:
            return 'rd'
        case _:
            return 'th'


def tier_name_to_color(tier: str):
    tier_names = ["F", "D-", "D", "D+", "C-", "C",
                  "C+", "B-", "B", "B+", "A-", "A", "A+",
                  "S", "S+"]
    tier_colors = ["00b050", "24bb45", "4bc735", "6fd226", "93de15", "b8e900",
                   "dcf400", "ffff00", "ffd600", "ffac00", "ff5d00", "ff5700",
                   "ff2b00", "ff0000", "ff0050"]

    if tier in tier_names:
        return RGBColor.from_string(tier_colors[tier_names.index(tier)])
    else:
        return RGBColor.from_string("ffffff")


def update_text(shape, text: str):
    for paragraph in shape.text_frame.paragraphs:
        for child in paragraph.runs:
            if child.text and child.text.strip():
                child.text = text


def replace_text(shape, text_map: dict[str, str]):
    for paragraph in shape.text_frame.paragraphs:
        for child in paragraph.runs:
            # print(repr(child.text))
            if child.text in text_map:
                child.text = text_map[child.text]


def replace_table_text(shape, text_map: dict[str, str]):
    for cell in shape.table.iter_cells():
        if cell.text in text_map:
            update_text(cell, text_map[cell.text])


def swap_image(shape, new_image: str):
    image, rid = shape.part.get_or_add_image_part(new_image)
    shape._pic.blipFill.blip.rEmbed = rid


def get_move_type(move_id):
    return moves[move_id]["type"]


def desaturate_color(color, amount):
    r, g, b = color
    h, s, v = colorsys.rgb_to_hsv(r / 255, g / 255, b / 255)
    new_r, new_g, new_b = colorsys.hsv_to_rgb(h, s * amount, 1.0)
    return RGBColor(int(new_r * 255), int(new_g * 255), int(new_b * 255))


def populate_move(child, trainer, mon_index, move_index):
    move_data = [moveset for moveset in raw_trainer_movesets if
                 moveset["class"] == trainer.class_id and moveset["instance"] == trainer.instance_id][0]
    has_move = mon_index < len(trainer.pokémon) and (move_id := move_data["moves"][mon_index][move_index])
    if has_move:
        update_text(child, format_name(MoveIdentifier.Name(move_id)))

        # Set move type color
        move_type = get_move_type(move_id)
        child.fill.fore_color.rgb = type_colors[move_type]

        # Set text to white if type has dark background
        if move_type in ["FIGHTING", "POISON", "GHOST", "DARK", "DRAGON"]:
            child.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string("ffffff")
    else:
        update_text(child, "")

        # set to light gray color
        child.fill.fore_color.rgb = RGBColor.from_string("d3d3d3")


def populate_mon(shape, trainer, mon_index):
    has_mon = mon_index < len(trainer.pokémon)
    if has_mon:
        species_type = species_types[trainer.pokémon[mon_index].species]
    for child in shape.shapes:
        match child.name:
            case "move-1" | "move-2" | "move-3" | "move-4":
                move_index = int(child.name[-1]) - 1
                populate_move(child, trainer, mon_index, move_index)
            case "held-item":
                if has_mon and trainer.pokémon[mon_index].held_item:
                    swap_image(child,
                               f"./trainer_cards/item_sprites/scaled/{item_sprites[trainer.pokémon[mon_index].held_item]}")
                else:
                    swap_image(child, "./trainer_cards/nothing_image.png")
            case "name-and-level":
                if has_mon:
                    replace_text(child, {
                        "Magikarp": format_name(trainer.pokémon[mon_index].species),
                        "Lv. 18": f"Lv. {trainer.pokémon[mon_index].level} ",
                    })

                else:
                    update_text(child, "")
            case "sprite":
                if has_mon:
                    dex_number = pokemon_in_order.index(format_name(trainer.pokémon[mon_index].species)) + 1
                    swap_image(child, f"./trainer_cards/pokemon_sprites/scaled/{dex_number:03}MS6.png")
                else:
                    swap_image(child, "./trainer_cards/pokemon_sprites/scaled/XYLoadingMS-2.png")
            case "sprite-backdrop":
                pass
            case "card-backdrop":
                if has_mon:
                    # Set gradient based on types
                    child.fill.gradient_stops[0].color.rgb = desaturate_color(type_colors[species_type[0]], 0.5)
                    r, g, b = type_colors[species_type[1]]
                    h, s, v = colorsys.rgb_to_hsv(r / 255, g / 255, b / 255)
                    new_r, new_g, new_b = colorsys.hsv_to_rgb(h, s * 0.5, v)
                    child.fill.gradient_stops[1].color.rgb = desaturate_color(type_colors[species_type[1]], 0.5)
                else:
                    # set to light gray color
                    child.fill.solid()
                    child.fill.fore_color.rgb = RGBColor.from_string("d3d3d3")


def format_trainer_name(trainer: Trainer):
    suffix =  f"#{trainer.rematch}" if trainer.rematch != 1 or trainer.has_later_rematch else ""
    return f"{trainer.class_name} {trainer.name} {suffix}"


def main():
    prs = Presentation("./trainer_cards/trainer_card_base.pptm")
    print("loading augmented trainer data")
    with open("omega_augmented_trainer_list.pickle", "rb") as f:
        augmented_trainer_list: List[Trainer] = pickle.loads(zlib.decompress(f.read()))
    print("loaded")

    augmented_trainer_list.reverse()

    item_set = set()

    for trainer in augmented_trainer_list:
        for mon in trainer.pokémon:
            item_set.add(mon.held_item)

        trainer_items = trainer.items
        for item in trainer_items:
            item_set.add(item)

    pprint(list(item_set))

    win_record_bar_width = 7367411

    print("updating slides")
    for i, slide in tqdm(enumerate(prs.slides), total=541, unit="slides"):
        trainer = augmented_trainer_list[i]
        if trainer.is_unused:
            trainer_index = 0
        elif trainer.is_rematch:
            original_battle = [ot for ot in augmented_trainer_list if
                               not ot.is_rematch and ot.class_id == trainer.class_id and ot.name == trainer.name][0]
            trainer_index = original_battle.game_index - 1
        else:
            trainer_index = trainer.game_index - 1

        trainer.victories.sort(key=lambda b: b.losing_trainer.elo, reverse=True)
        greatest_victory = trainer.victories[0]
        greatest_victory_trainer = greatest_victory.losing_trainer

        trainer.defeats.sort(key=lambda b: b.winning_trainer.elo)
        worst_defeat = trainer.defeats[0]
        worst_defeat_trainer = worst_defeat.winning_trainer

        for shape in slide.shapes:
            match shape.name:
                case "party-mon-1" | "party-mon-2" | "party-mon-3" | "party-mon-4" | "party-mon-5" | "party-mon-6":
                    mon_index = int(shape.name[-1]) - 1
                    populate_mon(shape, trainer, mon_index)
                case "strats":
                    for child in shape.shapes:
                        strat = child.name.split("-")[1]
                        trainer_strats = ["none" if strat == "NO_AI" else strat.split("_")[1].lower() for strat in trainer.strategy]
                        if strat not in trainer_strats:
                            child.fill.fore_color.rgb = desaturate_color(child.fill.fore_color.rgb, 0.2)
                            #update text color
                            child.text_frame.paragraphs[0].runs[0].font.color.rgb = desaturate_color(child.fill.fore_color.rgb, 1.2)
                        # swap cautious and setup?
                        # or pair them logically?
                        # smart - none - basic
                        # cautious - risky
                        # types - status
                        # offensive - opportunist
                        # setup - aggressive
                case "map-pin":
                    if trainer.is_unused:
                        swap_image(shape, "./trainer_cards/nothing_image.png")
                        continue

                    offset = int(trainer_index / 395 * win_record_bar_width)
                    shape.left = shape.left + offset

                    if trainer.is_rematch:
                        if trainer.continent == "Kanto":
                            swap_image(shape, "./trainer_cards/poke-ball-transparent.png")
                        else:
                            swap_image(shape, "./trainer_cards/gs-ball-transparent.png")
                    else:
                        if trainer.continent == "Kanto":
                            swap_image(shape, "./trainer_cards/pixel-pokeball-pin-large.png")

                case "w-l-ratio":
                    replace_table_text(shape, {
                        "1523": str(len(trainer.victories)),
                        "1719": str(len(trainer.defeats))
                    })
                case "dv-table":
                    replace_table_text(shape, {
                        "15": str(trainer.dvs[0]),
                        "17": str(trainer.dvs[1]),
                        "12": str(trainer.dvs[2]),
                        "19": str(trainer.dvs[3]),
                    })
                case "held-item-1" | "held-item-2":
                    item_number = int(shape.name[-1]) - 1
                    if item_number < len(trainer.items):
                        item = trainer.items[item_number]
                        swap_image(shape, f"./trainer_cards/item_sprites/scaled/{item_sprites[item]}")
                    else:
                        swap_image(shape, "./trainer_cards/item_sprites/scaled/no_item.png")
                case "items-label":
                    pass
                case "rank":
                    replace_text(shape, {
                        "#391": f"#{trainer.rank - 1}"
                    })
                case "ai-label":
                    pass
                case "dv-label":
                    pass
                case "gender-symbol":
                    if trainer.gender_symbol != "♂":
                        replace_text(shape, {
                            "♂": ""
                        })
                    if trainer.gender_symbol != "♀":
                        replace_text(shape, {
                            "♀": ""
                        })
                    if trainer.gender_symbol != "⚥":
                        replace_text(shape, {
                            "⚥": ""
                        })
                case "game-location":
                    if not trainer.is_unused and not trainer.is_rematch:
                        location = trainer.area
                        order = trainer.game_index
                        ordinal = getSuffix(order)
                        replace_text(shape, {
                            "Cherrygrove City – 3": f"{location} - {order}",
                            "rd ": ordinal + " ",
                            "overall": "in-game battle"
                        })
                    elif not trainer.is_unused and trainer.is_rematch:
                        location = trainer.area
                        rematch_count = trainer.rematch - 1
                        ordinal = getSuffix(rematch_count)
                        replace_text(shape, {
                            "Cherrygrove City – 3": f"{location} - ",
                            "rd ": "",
                            "overall": "rematch"
                        })
                    else:
                        replace_text(shape, {
                            "Cherrygrove City – 3": "Unused trainer",
                            "rd ": "",
                            "overall": ""
                        })

                    if trainer_index > 395 / 2:
                        shape.left = shape.left - shape.width - 250000
                        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    offset = int(trainer_index / 395 * win_record_bar_width)
                    shape.left = shape.left + offset
                case "win-record":
                    swap_image(shape, f"./trainer_cards/win_records/{trainer.class_id}-{trainer.instance_id}.png")
                case "greatest-win-pic":
                    swap_image(shape, f"../video-crystal/trainer_sprites/large/{trainer_sprites[greatest_victory_trainer.class_id]}")

                    shape.shadow._element.getchildren()[3].getchildren()[0].getchildren()[0].attrib["val"] = str(tier_name_to_color(greatest_victory_trainer.tier))
                case "greatest-win-box":
                    has_suffix = greatest_victory_trainer.rematch != 1 or greatest_victory_trainer.has_later_rematch
                    replace_text(shape, {
                        "C-": greatest_victory_trainer.tier,
                        " Hiker MICHELLE ": f" {format_name(greatest_victory_trainer.class_name)} {format_name(greatest_victory_trainer.name)} ",
                        "#": "#" if has_suffix else "",
                        '14 (1244)': f"{greatest_victory_trainer.rematch} ({round(greatest_victory_trainer.elo)})" if has_suffix else f"({round(greatest_victory_trainer.elo)})",
                        "mxrb ji5x fwty": prettify_hashid(greatest_victory.seed)
                    })

                    shape.text_frame.paragraphs[1].runs[0].font.color.rgb = tier_name_to_color(greatest_victory_trainer.tier)
                case "worst-loss-pic":
                    swap_image(shape, f"../video-crystal/trainer_sprites/large/{trainer_sprites[worst_defeat_trainer.class_id]}")

                    shape.shadow._element.getchildren()[3].getchildren()[0].getchildren()[0].attrib["val"] = str(tier_name_to_color(worst_defeat_trainer.tier))
                case "worst-loss-box":
                    has_suffix = worst_defeat_trainer.rematch != 1 or worst_defeat_trainer.has_later_rematch
                    replace_text(shape, {
                        "F": worst_defeat_trainer.tier,
                        " Fisher MICHELLE ": f" {format_name(worst_defeat_trainer.class_name)} {format_name(worst_defeat_trainer.name)} ",
                        "#": "#" if has_suffix else "",
                        '11 (443)': f"{worst_defeat_trainer.rematch} ({round(worst_defeat_trainer.elo)})" if has_suffix else f"({round(worst_defeat_trainer.elo)})",
                        "3q5t xto! trtx": prettify_hashid(worst_defeat.seed)
                    })

                    shape.text_frame.paragraphs[1].runs[0].font.color.rgb = tier_name_to_color(worst_defeat_trainer.tier)
                case "elo-label": pass
                case "elo":
                    replace_text(shape, {
                        "2347": str(round(trainer.elo))
                    })
                case "win-bar":
                    shape.width = int(3829811 * len(trainer.victories) / len(trainer.battles))
                case "lose-bar": pass
                case "name":
                    replace_text(shape, {
                        "Rival": format_name(trainer.class_name),
                        "??? ": format_name(trainer.name) + " ",
                        "#27": f"#{trainer.rematch}" if trainer.rematch != 1 or trainer.has_later_rematch else ""
                    })
                case "tier":
                    update_text(shape, trainer.tier)
                    shape.fill.fore_color.rgb = tier_name_to_color(trainer.tier)
                case "map-backdrop":
                    if trainer.is_unused:
                        swap_image(shape, "./trainer_cards/unused_grass.png")
                case "card-bg": pass
                case "trainer-sprite":
                    swap_image(shape, f"../video-crystal/trainer_sprites/large/{trainer_sprites[trainer.class_id]}")
                    shape.shadow._element.getchildren()[2].getchildren()[0].getchildren()[0].attrib["val"] = str(tier_name_to_color(trainer.tier))

    print("saving presentation")
    prs.save("./trainer_cards/trainer_card_test.pptm")


if __name__ == '__main__':
    with open("./analysis_results/most_used_moves.txt", 'r') as f:
        lines = f.readlines()

    for line in lines:
        move, count = line.split()
        print(format_name(move), "\t", count)



    main()
